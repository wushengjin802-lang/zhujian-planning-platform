from __future__ import annotations

import io
import time
import zipfile

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from sqlalchemy import func, select, text
from sqlalchemy.orm import Session

from app.db.models import Artifact, ChapterCitation, FactItem, Project, QualityIssue, ReportChapter
from app.services.office_converter import convert_bytes_to_pdf
from app.services.storage import storage


def generate_artifact(db: Session, artifact_id: str) -> Artifact | None:
    artifact = db.get(Artifact, artifact_id)
    if not artifact:
        return None
    blockers = db.scalar(
        select(func.count())
        .select_from(QualityIssue)
        .where(QualityIssue.project_id == artifact.project_id, QualityIssue.severity == "阻断", QualityIssue.status != "已关闭")
    )
    if blockers:
        artifact.status = "受阻"
        artifact.updated_at = current_day(db)
        db.commit()
        return artifact

    project = db.get(Project, artifact.project_id)
    chapters = db.scalars(select(ReportChapter).where(ReportChapter.project_id == artifact.project_id).order_by(ReportChapter.chapter_no)).all()
    facts = db.scalars(select(FactItem).where(FactItem.project_id == artifact.project_id)).all()
    issues = db.scalars(select(QualityIssue).where(QualityIssue.project_id == artifact.project_id)).all()
    citations = db.scalars(select(ChapterCitation).where(ChapterCitation.chapter_id.in_([chapter.id for chapter in chapters])).order_by(ChapterCitation.chapter_id, ChapterCitation.id)).all() if chapters else []
    if artifact.format == "Excel":
        data, mime = build_excel(project, chapters, facts, issues, citations)
    elif artifact.format == "Archive":
        data, mime = build_archive(project, chapters, facts, issues, citations)
    elif artifact.format == "PPT":
        data, mime = build_ppt(project, chapters, facts, issues)
    else:
        data, mime = build_word(project, chapters, facts, issues, citations)

    safe_name = artifact.name.replace("/", "_").replace("\\", "_")
    key = f"artifacts/{int(time.time() * 1000)}-{safe_name}"
    artifact.storage_path, artifact.file_size = storage.put_bytes(data, key, mime)
    artifact.status = "已生成"
    artifact.generated_at = db.scalar(select(func.now()))
    artifact.updated_at = current_day(db)
    db.commit()
    return artifact


def build_word(project: Project | None, chapters: list[ReportChapter], facts: list[FactItem], issues: list[QualityIssue], citations: list[ChapterCitation]) -> tuple[bytes, str]:
    doc = Document()
    doc.add_heading(project.name if project else "项目成果报告", level=1)
    doc.add_paragraph("本成果由平台自动生成。正式发布前需完成专业复核、质量门禁关闭和项目负责人确认。")
    doc.add_heading("一、项目基础信息", level=2)
    if project:
        table = doc.add_table(rows=0, cols=2)
        for label, value in [
            ("项目类型", project.type),
            ("项目地点", project.location),
            ("当前阶段", project.phase),
            ("项目负责人", project.owner),
            ("总体进度", f"{project.progress}%"),
            ("风险等级", project.risk),
        ]:
            cells = table.add_row().cells
            cells[0].text = label
            cells[1].text = value

    doc.add_heading("二、事实与指标底板", level=2)
    fact_table = doc.add_table(rows=1, cols=5)
    for index, label in enumerate(["分组", "指标", "取值", "状态", "来源"]):
        fact_table.rows[0].cells[index].text = label
    for fact in facts:
        cells = fact_table.add_row().cells
        cells[0].text = fact.fact_group
        cells[1].text = fact.name
        cells[2].text = f"{fact.value}{fact.unit or ''}"
        cells[3].text = fact.status
        cells[4].text = fact.source

    doc.add_heading("三、报告章节", level=2)
    for chapter in chapters:
        doc.add_heading(f"{chapter.chapter_no}. {chapter.title}", level=3)
        doc.add_paragraph(chapter.content or f"章节状态：{chapter.status}；引用 {chapter.citation_count} 条。")
        chapter_citations = [citation for citation in citations if citation.chapter_id == chapter.id]
        if chapter_citations:
            doc.add_paragraph("引用依据：")
            for citation in chapter_citations:
                doc.add_paragraph(f"{citation.excerpt}（{citation.source}）", style="List Bullet")

    doc.add_heading("四、质量门禁", level=2)
    for issue in issues:
        doc.add_paragraph(f"[{issue.severity}] {issue.title} - {issue.status}", style="List Bullet")
    stream = io.BytesIO()
    doc.save(stream)
    return stream.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def build_excel(project: Project | None, chapters: list[ReportChapter], facts: list[FactItem], issues: list[QualityIssue], citations: list[ChapterCitation]) -> tuple[bytes, str]:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "投资与事实底板"
    sheet.append(["项目", project.name if project else ""])
    sheet.append([])
    sheet.append(["分组", "事实项", "值", "单位", "来源", "状态"])
    for fact in facts:
        sheet.append([fact.fact_group, fact.name, fact.value, fact.unit, fact.source, fact.status])
    chapter_sheet = workbook.create_sheet("章节与引用")
    chapter_sheet.append(["章节号", "章节标题", "状态", "引用数", "引用摘录", "来源"])
    for chapter in chapters:
        chapter_citations = [citation for citation in citations if citation.chapter_id == chapter.id]
        if not chapter_citations:
            chapter_sheet.append([chapter.chapter_no, chapter.title, chapter.status, chapter.citation_count, "", ""])
        for citation in chapter_citations:
            chapter_sheet.append([chapter.chapter_no, chapter.title, chapter.status, chapter.citation_count, citation.excerpt, citation.source])
    issue_sheet = workbook.create_sheet("质量问题")
    issue_sheet.append(["等级", "类型", "问题", "责任人", "状态"])
    for issue in issues:
        issue_sheet.append([issue.severity, issue.type, issue.title, issue.owner, issue.status])
    stream = io.BytesIO()
    workbook.save(stream)
    return stream.getvalue(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


def build_archive(project: Project | None, chapters: list[ReportChapter], facts: list[FactItem], issues: list[QualityIssue], citations: list[ChapterCitation]) -> tuple[bytes, str]:
    word, _ = build_word(project, chapters, facts, issues, citations)
    excel, _ = build_excel(project, chapters, facts, issues, citations)
    ppt, _ = build_ppt(project, chapters, facts, issues)
    pdf = convert_bytes_to_pdf(word, "成果报告.docx")
    stream = io.BytesIO()
    with zipfile.ZipFile(stream, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("成果报告.docx", word)
        archive.writestr("事实与质量清单.xlsx", excel)
        archive.writestr("汇报简版.pptx", ppt)
        archive.writestr("决策记录.txt", build_decision_record(project, chapters, facts, issues, citations))
        if pdf:
            archive.writestr("成果报告.pdf", pdf)
    return stream.getvalue(), "application/zip"


def build_decision_record(project: Project | None, chapters: list[ReportChapter], facts: list[FactItem], issues: list[QualityIssue], citations: list[ChapterCitation]) -> str:
    lines = [
        f"项目：{project.name if project else '项目成果'}",
        "阶段：项目成果归档",
        f"章节数：{len(chapters)}",
        f"事实指标数：{len(facts)}",
        f"引用数：{len(citations)}",
        f"未关闭质量问题数：{sum(1 for issue in issues if issue.status != '已关闭')}",
        "",
        "发布判断：所有阻断问题关闭、关键章节审核通过后方可正式发布。",
    ]
    return "\n".join(lines)


def build_ppt(project: Project | None, chapters: list[ReportChapter], facts: list[FactItem], issues: list[QualityIssue]) -> tuple[bytes, str]:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = project.name if project else "项目成果汇报"
    slide.placeholders[1].text = "由平台自动生成的阶段性汇报材料，需专业人员复核。"

    summary = presentation.slides.add_slide(presentation.slide_layouts[1])
    summary.shapes.title.text = "事实与指标摘要"
    body = summary.placeholders[1].text_frame
    body.clear()
    for fact in facts[:6]:
        paragraph = body.add_paragraph()
        paragraph.text = f"{fact.name}：{fact.value}{fact.unit or ''}"
        paragraph.level = 0

    chapter_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    chapter_slide.shapes.title.text = "报告章节进度"
    chapter_body = chapter_slide.placeholders[1].text_frame
    chapter_body.clear()
    for chapter in chapters[:8]:
        paragraph = chapter_body.add_paragraph()
        paragraph.text = f"{chapter.chapter_no}. {chapter.title} - {chapter.status}"
        paragraph.level = 0

    issue_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    issue_slide.shapes.title.text = "质量门禁"
    issue_body = issue_slide.placeholders[1].text_frame
    issue_body.clear()
    for issue in issues[:6]:
        paragraph = issue_body.add_paragraph()
        paragraph.text = f"[{issue.severity}] {issue.title} - {issue.status}"
        paragraph.level = 0

    stream = io.BytesIO()
    presentation.save(stream)
    return stream.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def current_day(db: Session) -> str:
    return db.scalar(select(text("to_char(current_date, 'YYYY-MM-DD')")))
